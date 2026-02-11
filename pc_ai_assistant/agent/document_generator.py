"""
Document Generator using Python libraries
Creates documents locally and provides download links
"""
import os
from datetime import datetime
from pathlib import Path


def create_word_document(title, content, output_dir="documents"):
    """
    Create a Word document using python-docx
    
    Args:
        title: Document title
        content: Content (string or list of paragraphs)
        output_dir: Directory to save the document
    
    Returns:
        Path to the created document
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("⚠️ python-docx not installed. Installing...")
        import subprocess
        subprocess.check_call(["pip", "install", "python-docx"])
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    print(f"Creating Word document: {title}")
    
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Create document
    doc = Document()
    
    # Add title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Add timestamp
    timestamp = datetime.now().strftime("%B %d, %Y at %I:%M %p")
    date_para = doc.add_paragraph(f"Created: {timestamp}")
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Add spacing
    doc.add_paragraph()
    
    # Add content
    if isinstance(content, list):
        for paragraph in content:
            if paragraph.strip():
                p = doc.add_paragraph(paragraph)
                p.style = 'Normal'
    else:
        doc.add_paragraph(content)
    
    # Save document
    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = os.path.join(output_dir, filename)
    doc.save(filepath)
    
    print(f"✅ Word document created: {filepath}")
    return filepath


def create_pdf_document(title, content, output_dir="documents"):
    """
    Create a PDF document using reportlab
    
    Args:
        title: Document title
        content: Content (string or list of paragraphs)
        output_dir: Directory to save the document
    
    Returns:
        Path to the created document
    """
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    except ImportError:
        print("⚠️ reportlab not installed. Installing...")
        import subprocess
        subprocess.check_call(["pip", "install", "reportlab"])
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    
    print(f"Creating PDF document: {title}")
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Create filename
    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    filepath = os.path.join(output_dir, filename)
    
    # Create PDF
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor='#1a1a2e',
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    date_style = ParagraphStyle(
        'DateStyle',
        parent=styles['Normal'],
        fontSize=10,
        textColor='#666666',
        spaceAfter=20,
        alignment=TA_CENTER
    )
    
    body_style = ParagraphStyle(
        'BodyStyle',
        parent=styles['Normal'],
        fontSize=12,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=12
    )
    
    # Add title
    story.append(Paragraph(title, title_style))
    
    # Add timestamp
    timestamp = datetime.now().strftime("%B %d, %Y at %I:%M %p")
    story.append(Paragraph(f"Created: {timestamp}", date_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Add content
    if isinstance(content, list):
        for paragraph in content:
            if paragraph.strip():
                story.append(Paragraph(paragraph, body_style))
    else:
        for paragraph in content.split('\n'):
            if paragraph.strip():
                story.append(Paragraph(paragraph, body_style))
    
    # Build PDF
    doc.build(story)
    
    print(f"✅ PDF document created: {filepath}")
    return filepath


def create_powerpoint(title, slides_content, output_dir="documents"):
    """
    Create a PowerPoint presentation using python-pptx
    
    Args:
        title: Presentation title
        slides_content: List of dicts with 'title' and 'content' for each slide
        output_dir: Directory to save the presentation
    
    Returns:
        Path to the created presentation
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("⚠️ python-pptx not installed. Installing...")
        import subprocess
        subprocess.check_call(["pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    
    print(f"Creating PowerPoint presentation: {title}")
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title_shape.text = title
    timestamp = datetime.now().strftime("%B %d, %Y")
    subtitle.text = f"Created by PC AI Assistant\n{timestamp}"
    
    # Content slides
    for slide_data in slides_content:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_data.get('title', 'Slide')
        
        tf = body_shape.text_frame
        content = slide_data.get('content', [])
        
        if isinstance(content, list):
            for i, point in enumerate(content):
                if i == 0:
                    tf.text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
        else:
            tf.text = content
    
    # Save presentation
    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    
    print(f"✅ PowerPoint presentation created: {filepath}")
    return filepath


def create_markdown_document(title, content, output_dir="documents"):
    """
    Create a Markdown document
    
    Args:
        title: Document title
        content: Content (string or list of paragraphs)
        output_dir: Directory to save the document
    
    Returns:
        Path to the created document
    """
    print(f"Creating Markdown document: {title}")
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Create markdown content
    md_content = f"# {title}\n\n"
    md_content += f"*Created: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}*\n\n"
    md_content += "---\n\n"
    
    if isinstance(content, list):
        for paragraph in content:
            if paragraph.strip():
                md_content += f"{paragraph}\n\n"
    else:
        md_content += content
    
    # Save document
    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
    filepath = os.path.join(output_dir, filename)
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    print(f"✅ Markdown document created: {filepath}")
    return filepath


def create_html_document(title, content, output_dir="documents"):
    """
    Create an HTML document
    
    Args:
        title: Document title
        content: Content (string or list of paragraphs)
        output_dir: Directory to save the document
    
    Returns:
        Path to the created document
    """
    print(f"Creating HTML document: {title}")
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Create HTML content
    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Arial, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 40px 20px;
            line-height: 1.6;
            color: #333;
        }}
        h1 {{
            color: #1a1a2e;
            border-bottom: 3px solid #667eea;
            padding-bottom: 10px;
        }}
        .meta {{
            color: #666;
            font-size: 14px;
            margin-bottom: 30px;
        }}
        p {{
            margin-bottom: 15px;
            text-align: justify;
        }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    <div class="meta">Created: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</div>
"""
    
    if isinstance(content, list):
        for paragraph in content:
            if paragraph.strip():
                html_content += f"    <p>{paragraph}</p>\n"
    else:
        for paragraph in content.split('\n'):
            if paragraph.strip():
                html_content += f"    <p>{paragraph}</p>\n"
    
    html_content += """</body>
</html>"""
    
    # Save document
    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
    filepath = os.path.join(output_dir, filename)
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"✅ HTML document created: {filepath}")
    return filepath


# Example usage
if __name__ == "__main__":
    sample_title = "Sample Document"
    sample_content = [
        "Introduction",
        "This is a sample document created by the PC AI Assistant.",
        "Main Content",
        "Here is the main content of the document with multiple paragraphs.",
        "Conclusion",
        "This document was created automatically using Python libraries."
    ]
    
    # Test document creation
    print("Testing document generation...")
    word_path = create_word_document(sample_title, sample_content)
    pdf_path = create_pdf_document(sample_title, sample_content)
    md_path = create_markdown_document(sample_title, sample_content)
    html_path = create_html_document(sample_title, sample_content)
    
    print("\n✅ All documents created successfully!")
    print(f"Word: {word_path}")
    print(f"PDF: {pdf_path}")
    print(f"Markdown: {md_path}")
    print(f"HTML: {html_path}")
